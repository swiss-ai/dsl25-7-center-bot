import io
import logging
import zipfile
from googleapiclient.http import MediaIoBaseDownload
from pptx import Presentation

logger = logging.getLogger(__name__)

def extract_text_from_pptx(service, file_id):
    request = service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while not done:
        status, done = downloader.next_chunk()

    fh.seek(0)

    try:
        prs = Presentation(fh)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except zipfile.BadZipFile:
        logger.warning(f"❌ Failed to parse PPTX file (not a real zip): {file_id}")
        return ""
