import os
import json
import logging
import asyncio
import tempfile
import uuid
import pickle
from typing import Dict, List, Any, Optional, Tuple, Set
from datetime import datetime
from io import BytesIO

# Google API imports
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# Local imports
from services.knowledge.document_processor import DocumentProcessor
from db.vector_db import VectorDB

logger = logging.getLogger(__name__)

# Define the scopes
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

# Supported file types and their extractors
SUPPORTED_MIME_TYPES = {
    'application/pdf': 'extract_text_from_pdf',
    'text/plain': 'extract_text_from_txt',
    'text/markdown': 'extract_text_from_md',
    'application/vnd.google-apps.document': 'extract_text_from_docs',
    'application/vnd.google-apps.spreadsheet': 'extract_text_from_sheets',
    'application/vnd.google-apps.presentation': 'extract_text_from_slides',
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'extract_text_from_docx',
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'extract_text_from_xlsx',
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'extract_text_from_pptx',
}

# Provide a compatibility alias for the old GoogleDriveMCP class
class GoogleDriveMCP:
    """
    Legacy Google Drive MCP class for compatibility.
    This is just a placeholder to avoid import errors.
    """
    
    def __init__(self, credentials_path=None, token_path=None, document_processor=None):
        self.credentials_path = credentials_path
        self.token_path = token_path
        self.document_processor = document_processor
        self.service = None
        
        logger.warning("Legacy GoogleDriveMCP is deprecated. Please use GoogleDriveManager instead.")


class GoogleDriveManager:
    """
    Google Drive content management for the vector database.
    Based on the gdrive_scraping implementation.
    """
    
    def __init__(
        self, 
        document_processor: Optional[DocumentProcessor] = None,
        vector_db: Optional[VectorDB] = None,
        credentials_path: str = None, 
        token_path: str = None,
        sync_file: str = "gdrive_last_sync.json"
    ):
        """
        Initialize the Google Drive manager.
        
        Args:
            document_processor: DocumentProcessor instance for processing documents
            vector_db: VectorDB instance for storing documents
            credentials_path: Path to the credentials.json file
            token_path: Path to token.pickle file for authentication
            sync_file: Path to the file used to track synced files
        """
        self.document_processor = document_processor
        self.vector_db = vector_db
        self.credentials_path = credentials_path or os.getenv('GOOGLE_CREDENTIALS_PATH')
        self.token_path = token_path or os.getenv('GOOGLE_TOKEN_PATH')
        self.sync_file = sync_file
        self.service = None
        
        # If credential paths provided, authenticate
        if self.credentials_path and os.path.exists(self.credentials_path):
            self._authenticate()
    
    def _authenticate(self):
        """Authenticate with Google Drive API."""
        creds = None
        
        # Load saved credentials if they exist
        if os.path.exists(self.token_path):
            with open(self.token_path, 'rb') as token:
                creds = pickle.load(token)
        
        # If credentials don't exist or are invalid, get new ones
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(
                    self.credentials_path, SCOPES
                )
                creds = flow.run_local_server(port=0)
            
            # Save the credentials for the next run
            with open(self.token_path, 'wb') as token:
                pickle.dump(creds, token)
        
        # Build the Drive service
        self.service = build('drive', 'v3', credentials=creds)
        logger.info("Google Drive API authenticated successfully")
    
    def _load_synced_files(self) -> Dict[str, str]:
        """
        Load the list of already synced files.
        
        Returns:
            Dict mapping file IDs to modification times
        """
        if os.path.exists(self.sync_file):
            try:
                with open(self.sync_file, 'r') as f:
                    return json.load(f)
            except Exception as e:
                logger.error(f"Error loading synced files: {e}")
        
        return {}
    
    def _save_synced_files(self, synced_files: Dict[str, str]):
        """
        Save the list of synced files.
        
        Args:
            synced_files: Dict mapping file IDs to modification times
        """
        try:
            with open(self.sync_file, 'w') as f:
                json.dump(synced_files, f)
        except Exception as e:
            logger.error(f"Error saving synced files: {e}")
    
    def _mark_file_synced(self, file_id: str, modified_time: str):
        """
        Mark a file as synced by adding it to the synced files list.
        
        Args:
            file_id: Google Drive file ID
            modified_time: Modification time from Google Drive
        """
        synced = self._load_synced_files()
        synced[file_id] = modified_time
        self._save_synced_files(synced)
    
    async def list_files(
        self, 
        query: str = None, 
        page_size: int = 100, 
        order_by: str = "modifiedTime desc",
        fields: str = "files(id, name, mimeType, modifiedTime, createdTime, description, webViewLink)"
    ) -> List[Dict[str, Any]]:
        """
        List files in Google Drive.
        
        Args:
            query: Search query (Google Drive Query Language)
            page_size: Number of files to return
            order_by: Sort order
            fields: Fields to include in the response
            
        Returns:
            List of file metadata
        """
        if not self.service:
            logger.error("Google Drive API not authenticated")
            return []
        
        try:
            # Execute the request in a separate thread
            loop = asyncio.get_event_loop()
            
            results = []
            page_token = None
            
            while True:
                request = self.service.files().list(
                    q=query,
                    pageSize=page_size,
                    orderBy=order_by,
                    fields=f"nextPageToken, {fields}",
                    pageToken=page_token
                )
                
                response = await loop.run_in_executor(None, request.execute)
                results.extend(response.get('files', []))
                
                page_token = response.get('nextPageToken')
                if not page_token:
                    break
            
            return results
        
        except Exception as e:
            logger.error(f"Error listing files from Google Drive: {e}")
            return []
    
    async def search_files(self, query: str, max_results: int = 10) -> List[Dict[str, Any]]:
        """
        Search for files in Google Drive by name and content.
        
        Args:
            query: Search query
            max_results: Maximum number of results to return
            
        Returns:
            List of file metadata
        """
        # Format the query for Google Drive's query language
        # This searches in file name and full-text content
        formatted_query = f"fullText contains '{query}' or name contains '{query}'"
        
        # Add additional filter for supported MIME types
        mime_types_query = " or ".join([f"mimeType='{mime}'" for mime in SUPPORTED_MIME_TYPES.keys()])
        formatted_query = f"({formatted_query}) and ({mime_types_query}) and trashed=false"
        
        # Log the exact query being used
        logger.info(f"Searching Google Drive with query: {formatted_query}")
        
        results = await self.list_files(
            query=formatted_query,
            page_size=max_results
        )
        
        # Log what was found
        if results:
            logger.info(f"Found {len(results)} files in Google Drive matching '{query}'")
            for i, file in enumerate(results):
                logger.info(f"  {i+1}. {file.get('name', 'Untitled')} ({file.get('mimeType', 'Unknown type')})")
        else:
            logger.info(f"No files found in Google Drive matching '{query}'")
            
        return results
    
    async def get_files_to_update(self) -> List[Dict[str, Any]]:
        """
        Get a list of files that need to be updated in the vector database.
        
        Returns:
            List of file metadata for files that need to be updated
        """
        if not self.service:
            logger.error("Google Drive API not authenticated")
            return []
        
        # Load the list of already synced files
        synced_files = self._load_synced_files()
        
        # Build a query for supported file types
        mime_types_query = " or ".join([f"mimeType='{mime}'" for mime in SUPPORTED_MIME_TYPES.keys()])
        query = f"({mime_types_query}) and trashed=false"
        
        # Get the list of files
        files = await self.list_files(query=query)
        
        # Filter for files that need to be updated
        to_update = []
        for file in files:
            file_id = file['id']
            if file_id not in synced_files or file['modifiedTime'] > synced_files[file_id]:
                to_update.append(file)
        
        logger.info(f"Found {len(to_update)} files to update out of {len(files)} total files")
        return to_update
    
    async def extract_text_from_pdf(self, file_id: str) -> str:
        """
        Extract text from a PDF file.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Download the file
            request = self.service.files().get_media(fileId=file_id)
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            file_data.seek(0)
            
            # Use PyMuPDF to extract text
            import fitz  # PyMuPDF
            text = ""
            
            with fitz.open(stream=file_data.getvalue(), filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()
            
            return text
        
        except Exception as e:
            logger.error(f"Error extracting text from PDF (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_txt(self, file_id: str) -> str:
        """
        Extract text from a plain text file.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Download the file
            request = self.service.files().get_media(fileId=file_id)
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Convert to text
            return file_data.getvalue().decode('utf-8')
        
        except Exception as e:
            logger.error(f"Error extracting text from TXT (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_md(self, file_id: str) -> str:
        """
        Extract text from a Markdown file.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        # For Markdown, we'll use the same approach as TXT
        return await self.extract_text_from_txt(file_id)
    
    async def extract_text_from_docs(self, file_id: str) -> str:
        """
        Extract text from a Google Doc.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Export the file as plain text
            request = self.service.files().export_media(fileId=file_id, mimeType='text/plain')
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Convert to text
            return file_data.getvalue().decode('utf-8')
        
        except Exception as e:
            logger.error(f"Error extracting text from Google Doc (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_sheets(self, file_id: str) -> str:
        """
        Extract text from a Google Sheet.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Export the file as CSV
            request = self.service.files().export_media(fileId=file_id, mimeType='text/csv')
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Convert to text
            return file_data.getvalue().decode('utf-8')
        
        except Exception as e:
            logger.error(f"Error extracting text from Google Sheet (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_slides(self, file_id: str) -> str:
        """
        Extract text from Google Slides.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Export the file as plain text
            request = self.service.files().export_media(fileId=file_id, mimeType='text/plain')
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Convert to text
            return file_data.getvalue().decode('utf-8')
        
        except Exception as e:
            logger.error(f"Error extracting text from Google Slides (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_docx(self, file_id: str) -> str:
        """
        Extract text from a DOCX file.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Download the file
            request = self.service.files().get_media(fileId=file_id)
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Extract text using python-docx
            import docx
            document = docx.Document(file_data)
            text = "\n".join([p.text for p in document.paragraphs])
            return text
        
        except Exception as e:
            logger.error(f"Error extracting text from DOCX (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_xlsx(self, file_id: str) -> str:
        """
        Extract text from an XLSX file.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Download the file
            request = self.service.files().get_media(fileId=file_id)
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Extract text using pandas
            import pandas as pd
            
            text_parts = []
            xl = pd.ExcelFile(file_data)
            
            for sheet_name in xl.sheet_names:
                df = xl.parse(sheet_name)
                text_parts.append(f"Sheet: {sheet_name}")
                text_parts.append(df.to_string())
            
            return "\n\n".join(text_parts)
        
        except Exception as e:
            logger.error(f"Error extracting text from XLSX (ID: {file_id}): {e}")
            return ""
    
    async def extract_text_from_pptx(self, file_id: str) -> str:
        """
        Extract text from a PPTX file.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            Extracted text
        """
        if not self.service:
            return ""
        
        try:
            # Download the file
            request = self.service.files().get_media(fileId=file_id)
            file_data = BytesIO()
            downloader = MediaIoBaseDownload(file_data, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            # Extract text using python-pptx
            from pptx import Presentation
            
            presentation = Presentation(file_data)
            text_parts = []
            
            for i, slide in enumerate(presentation.slides):
                text_parts.append(f"Slide {i+1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            
            return "\n".join(text_parts)
        
        except Exception as e:
            logger.error(f"Error extracting text from PPTX (ID: {file_id}): {e}")
            return ""
    
    async def get_file_content(self, file_id: str, file_metadata: Dict[str, Any] = None) -> Tuple[str, Dict[str, Any]]:
        """
        Get the content of a file from Google Drive.
        
        Args:
            file_id: The ID of the file
            file_metadata: Optional file metadata to avoid an extra API call
            
        Returns:
            Tuple of (content, metadata)
        """
        if not self.service:
            logger.error("Google Drive API not authenticated")
            return "", {}
        
        try:
            # Get file metadata if not provided
            if not file_metadata:
                loop = asyncio.get_event_loop()
                file_request = self.service.files().get(
                    fileId=file_id, 
                    fields="id, name, mimeType, description, modifiedTime, createdTime, webViewLink, owners"
                )
                file_metadata = await loop.run_in_executor(None, file_request.execute)
            
            # Check if the MIME type is supported
            mime_type = file_metadata.get('mimeType')
            if mime_type not in SUPPORTED_MIME_TYPES:
                logger.warning(f"Unsupported MIME type: {mime_type}")
                return f"Unsupported file type: {mime_type}", file_metadata
            
            # Get the extractor method name
            extractor_name = SUPPORTED_MIME_TYPES[mime_type]
            extractor = getattr(self, extractor_name, None)
            
            if not extractor:
                logger.error(f"Extractor {extractor_name} not implemented")
                return "", file_metadata
            
            # Extract the text content
            content = await extractor(file_id)
            return content, file_metadata
        
        except Exception as e:
            logger.error(f"Error getting file content: {e}")
            return f"Error: {str(e)}", {}
    
    def _chunk_text(self, text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> List[str]:
        """
        Split text into chunks with overlap.
        
        Args:
            text: The text to split
            chunk_size: Maximum chunk size
            chunk_overlap: Overlap between chunks
            
        Returns:
            List of text chunks
        """
        if not text:
            return []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunks.append(text[start:end])
            start = end - chunk_overlap
        
        return chunks
    
    async def process_file_to_vector_db(self, file_id: str, file_metadata: Dict[str, Any] = None) -> bool:
        """
        Process a file and add it to the vector database.
        
        Args:
            file_id: The ID of the file
            file_metadata: Optional file metadata to avoid an extra API call
            
        Returns:
            bool: Success status
        """
        if not self.document_processor:
            logger.error("Document processor not available")
            return False
        
        try:
            # Get the file content and metadata
            content, metadata = await self.get_file_content(file_id, file_metadata)
            
            if not content:
                logger.error(f"Failed to get content for file {file_id}")
                return False
            
            # If the file already exists in the vector DB, remove it first
            # This ensures we don't have duplicate versions
            doc_id = f"gdrive_{file_id}"
            
            # Prepare metadata for the document processor
            doc_metadata = {
                "doc_id": doc_id,
                "title": metadata.get("name", "Untitled Google Drive Document"),
                "source": "google_drive",
                "url": metadata.get("webViewLink", ""),
                "author": self._get_authors(metadata),
                "created_at": metadata.get("createdTime", ""),
                "modified_at": metadata.get("modifiedTime", ""),
                "mime_type": metadata.get("mimeType", ""),
                "description": metadata.get("description", ""),
                "tags": "google_drive"  # ChromaDB expects string, not list
            }
            
            # Process the document
            chunk_ids = await self.document_processor.process_document(
                content=content,
                metadata=doc_metadata
            )
            
            # Mark the file as synced
            self._mark_file_synced(file_id, metadata.get("modifiedTime", ""))
            
            logger.info(f"Added Google Drive file {file_id} to vector database in {len(chunk_ids)} chunks")
            return True
        
        except Exception as e:
            logger.error(f"Error processing Google Drive file {file_id}: {e}")
            return False
    
    async def sync_all_files(self, max_files: int = None) -> int:
        """
        Sync all files from Google Drive to the vector database.
        
        Args:
            max_files: Maximum number of files to sync (None for all)
            
        Returns:
            int: Number of files synced
        """
        if not self.document_processor:
            logger.error("Document processor not available")
            return 0
        
        if not self.service:
            logger.error("Google Drive API not authenticated")
            return 0
        
        # Get files that need to be updated
        files_to_update = await self.get_files_to_update()
        
        if max_files:
            files_to_update = files_to_update[:max_files]
        
        if not files_to_update:
            logger.info("No files found to sync")
            return 0
        
        # Process each file
        synced_count = 0
        for file in files_to_update:
            success = await self.process_file_to_vector_db(file['id'], file)
            if success:
                synced_count += 1
        
        logger.info(f"Synced {synced_count} Google Drive files to vector database")
        return synced_count
    
    async def sync_recent_files(self, max_files: int = 20) -> int:
        """
        Sync only the most recently modified files.
        
        Args:
            max_files: Maximum number of files to sync
            
        Returns:
            int: Number of files synced
        """
        return await self.sync_all_files(max_files=max_files)
    
    def _get_authors(self, metadata: Dict[str, Any]) -> str:
        """Extract author information from file metadata."""
        owners = metadata.get("owners", [])
        if not owners:
            return "Unknown"
        
        names = [owner.get("displayName", "Unknown") for owner in owners]
        return ", ".join(names)